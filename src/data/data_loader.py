from __future__ import annotations

import numpy as np
import hashlib
import logging
from pathlib import Path
from typing import Iterable, List, Optional, Set

import fitz
from docx import Document as DocxDocument
from pptx import Presentation


from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_ollama import OllamaEmbeddings

from src import config

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}


def file_hash(file_path: str) -> str:
    """Generate a short content hash for change tracking."""
    try:
        return hashlib.md5(Path(file_path).read_bytes()).hexdigest()[:12]
    except FileNotFoundError:
        logger.warning(f"File not found for hashing: {file_path}")
        return ""


def _load_pdf(file_path: str) -> List[Document]:
    """Load PDF as one Document per page (metadata includes 1-based page number)."""

    p = Path(file_path)
    docs: List[Document] = []
    doc = fitz.open(str(p))
    try:
        base_meta = {
            "source": str(p),
            "doc_type": "pdf",
            "doc_hash": file_hash(str(p)),
        }
        temp = ""
        for page_index, page in enumerate(doc, start=1):
            text = (page.get_text("text") or "").strip()
            if not text or (len(text.split(" ")) < 20):
                temp += text
                continue
            if temp:
                text = temp + "\n" + text
                temp = ""
            meta = dict(base_meta)
            meta["page"] = page_index
            docs.append(Document(page_content=text, metadata=meta))
        return docs
    finally:
        doc.close()


def _paragraph_has_page_break(paragraph) -> bool:
    try:
        xml = paragraph._p.xml
    except Exception:
        return False
    return "w:type=\"page\"" in xml and "<w:br" in xml


def _load_docx(file_path: str) -> List[Document]:
    """Load DOCX as one Document per (approx) page.

    If the DOCX contains explicit page breaks, we split on those; otherwise this
    results in a single page=1 document.
    """
    p = Path(file_path)
    base_meta = {
        "source": str(p),
        "doc_type": "docx",
        "doc_hash": file_hash(str(p)),
    }

    doc = DocxDocument(str(p))
    docs: List[Document] = []
    current_page = 1
    buffer: List[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            buffer.append(text)

        if _paragraph_has_page_break(para):
            page_text = "\n".join([t for t in buffer if t.strip()]).strip()
            if page_text:
                meta = dict(base_meta)
                meta["page"] = current_page
                docs.append(Document(page_content=page_text, metadata=meta))
            buffer = []
            current_page += 1

    page_text = "\n".join([t for t in buffer if t.strip()]).strip()
    if page_text:
        meta = dict(base_meta)
        meta["page"] = current_page
        docs.append(Document(page_content=page_text, metadata=meta))

    return docs


def _load_pptx(file_path: str) -> List[Document]:
    """Load PPTX as one Document per slide (metadata includes 1-based slide number)."""

    p = Path(file_path)
    base_meta = {
        "source": str(p),
        "doc_type": "pptx",
        "doc_hash": file_hash(str(p)),
    }

    pres = Presentation(str(p))
    docs: List[Document] = []
    temp = ""
    for slide_index, slide in enumerate(pres.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            text = ""
            if hasattr(shape, "text"):
                try:
                    text = (shape.text or "").strip()
                except Exception:
                    text = ""
            if len(text.split(" ")) < 20:
                temp += text + "\n"
                continue
            if text:
                if temp:
                    parts.append(text + "\n" + temp)
                else:
                    parts.append(text)
        slide_text = "\n".join(parts).strip()
        if not slide_text:
            continue
        meta = dict(base_meta)
        meta["slide"] = slide_index
        docs.append(Document(page_content=slide_text, metadata=meta))
    return docs


def chunk_documents(
    docs: List[Document],
    method: str = "recursive",
) -> List[Document] | NotImplementedError:
    """Chunk documents with various strategies.

    Args:
        docs: List of Document objects to chunk.
        method: Chunking method - "recursive", "semantic", or "late". Default: "recursive"
    
    Returns:
        List of chunked Document objects.
    """
    if not docs:
        return []
    
    chunk_size = int(config.CHUNK_SIZE)
    chunk_overlap = int(config.CHUNK_OVERLAP)
    method = config.CHUNK_METHOD.lower().strip() if hasattr(config, "CHUNK_METHOD") else method.lower().strip()
    
    if method == "late":
        #TODO: Need to fix late chunking to get proper meta data
        return NotImplementedError("Late chunking not implemented yet.")
    
    if method == "semantic":
        try:
            splitter = RecursiveCharacterTextSplitter(chunk_size=300,
                                                      chunk_overlap=20)
            initial_chunks = splitter.split_documents(docs)
            
            if not initial_chunks:
                return docs
            
            embeddings = OllamaEmbeddings(
                base_url=config.OLLAMA_CONFIG["base_url"],
                model=config.EMBEDDING_MODEL_CONFIG["model"],
            )
            
            chunk_embeddings = embeddings.embed_documents([c.page_content for c in initial_chunks])
            
            merged = []
            current = initial_chunks[0]
            current_emb = np.array(chunk_embeddings[0])
            
            for i in range(1, len(initial_chunks)):
                next_chunk = initial_chunks[i]
                next_emb = np.array(chunk_embeddings[i])
                
                same_source = current.metadata.get("source") == next_chunk.metadata.get("source")
                similarity = np.dot(current_emb, next_emb) / (np.linalg.norm(current_emb) * np.linalg.norm(next_emb) + 1e-8)
                
                if same_source and similarity > 0.7:
                    merged_meta = current.metadata.copy()
                    for key in ["page", "slide"]:
                        if key in current.metadata or key in next_chunk.metadata:
                            vals = set()
                            if key in current.metadata:
                                vals.update([current.metadata[key]] if not isinstance(current.metadata[key], list) else current.metadata[key])
                            if key in next_chunk.metadata:
                                vals.update([next_chunk.metadata[key]] if not isinstance(next_chunk.metadata[key], list) else next_chunk.metadata[key])
                            merged_meta[key] = ",".join(sorted(list(vals)))
                    
                    current = Document(
                        page_content=current.page_content + "\n" + next_chunk.page_content,
                        metadata=merged_meta
                    )
                    current_emb = (current_emb + next_emb) / 2
                else:
                    merged.append(current)
                    current = next_chunk
                    current_emb = next_emb
            
            merged.append(current)
            return merged
            
        except Exception as e:
            logger.warning(f"Semantic chunking failed: {e}, falling back to recursive")
            method = "recursive"
    
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_documents(docs)

class DataLoader:
    """
    Main data loader that uses Docling.
    Provides a unified interface for loading documents from files or folders.
    """

    def __init__(self) -> None:
        pass

    @staticmethod
    def _normalize_extensions(extensions: Optional[Iterable[str]]) -> Set[str]:
        """
           Basically ensures extensions start with a dot and are lowercase.
        """
        if extensions is None:
            return SUPPORTED_EXTENSIONS
        
        normalized = set()
        for ext in extensions:
            ext = f".{ext.strip().lower().lstrip('.')}"
            if ext in SUPPORTED_EXTENSIONS:
                normalized.add(ext)
        return normalized

    def load(
        self,
        path: str,
        extensions: Optional[Iterable[str]] = None,
        recursive: bool = True,
        incremental: bool = False,
        existing_hashes: Optional[Set[str]] = None,
    ) -> List[Document]:
        """
        Load documents from a file or a folder.

        Args:
            path: The path to a single file or a directory.
            extensions: List of extensions to load (e.g., ["pdf", "docx"]). Defaults to all.
            recursive: If True, recursively load from subdirectories.
            incremental: If True, skip files whose hash is in existing_hashes.
            existing_hashes: A set of known file hashes to skip if incremental is True.

        Returns:
            A list of loaded Document objects.
        """
        p = Path(path)
        if not p.exists():
            logger.warning(f"Path does not exist: {path}")
            return []

        if p.is_file():
            if incremental and file_hash(str(p)) in (existing_hashes or set()):
                return []
            return self._load_file(str(p))

        if p.is_dir():
            return self._load_from_folder(p, extensions, recursive, incremental, existing_hashes)

        return []

    def _load_file(self, file_path: str) -> List[Document]:
        p = Path(file_path)
        ext = p.suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            logger.warning(f"Skipping unsupported file type: {file_path}")
            return []

        try:
            if ext == ".pdf":
                return _load_pdf(str(p))
            if ext == ".docx":
                return _load_docx(str(p))
            if ext == ".pptx":
                return _load_pptx(str(p))
        except Exception as e:
            logger.error(f"Failed to process file {file_path}: {e}", exc_info=True)

        return []

    def _load_from_folder(
        self,
        folder_path: Path,
        extensions: Optional[Iterable[str]],
        recursive: bool,
        incremental: bool,
        existing_hashes: Optional[Set[str]],
    ) -> List[Document]:
        docs: List[Document] = []
        valid_exts = self._normalize_extensions(extensions)
        seen_hashes = existing_hashes or set()

        iterator = folder_path.rglob("*") if recursive else folder_path.iterdir()
        
        for file_path in iterator:
            if not file_path.is_file() or file_path.suffix.lower() not in valid_exts:
                continue

            if incremental:
                h = file_hash(str(file_path))
                if h and h in seen_hashes:
                    continue
            
            docs.extend(self._load_file(str(file_path)))

        return docs


def load_and_chunk(path: str, *, recursive: bool = True) -> List[Document]:
    """Convenience API: load files and return chunked Documents."""

    loader = DataLoader()
    docs = loader.load(path, recursive=recursive)
    return chunk_documents(docs)